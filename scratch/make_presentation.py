"""
make_presentation.py — Generates publication-grade presentation slides for viva.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

DARK_BLUE = RGBColor(30, 58, 138)
ACCENT_BLUE = RGBColor(59, 130, 246)
TEXT_DARK = RGBColor(31, 41, 55)
BG_LIGHT = RGBColor(248, 250, 252)

def add_header(slide, title_text, subtitle_text=""):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(100, 116, 139)

# ── Slide 1: Title Slide ──────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
tf1 = tb1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "IPC2BNS-Verify"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

p2 = tf1.add_paragraph()
p2.text = "A Constraint-Verified, Incrementally Refreshable RAG Architecture for Indian Statutory Transitions"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT_BLUE

p3 = tf1.add_paragraph()
p3.text = "\nSubstantive Law (IPC 1860 → BNS 2023) & Procedural Law (CrPC 1973 → BNSS 2023)\nDepartment of Computer Science & Engineering | September 2026"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_DARK

# ── Slide 2: The Problem ──────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_header(slide2, "The Problem: Historical Inertia & Legal LLM Hallucinations", "Why standard RAG fails during major statutory transitions")
tb2 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
tf2 = tb2.text_frame
tf2.word_wrap = True

bullets2 = [
    "• 164 Years of Historical Pre-Training Bias: 99%+ of Indian legal texts cite IPC 1860 & CrPC 1973.",
    "• Historical Inertia in Closed-Book LLMs: Foundation models default to obsolete sections (10.0% accuracy on current law).",
    "• Force-Mapping Repealed Offences: Standard RAG force-maps repealed laws (Sedition §124A, Adultery §497) into wrong sections.",
    "• Right Section, Wrong Question: Models cite valid sections that are completely non-responsive to the query.",
    "• Cross-Statute Inconsistencies: Contradictory citations (citing IPC §302 Murder alongside BNS §318 Cheating)."
]
for b in bullets2:
    p = tf2.add_paragraph()
    p.text = b
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(12)

# ── Slide 3: System Architecture ──────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_header(slide3, "System Architecture: Decoupled Constraint Verification", "End-to-End Neuro-Symbolic Pipeline")
tb3 = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
tf3 = tb3.text_frame
tf3.word_wrap = True

bullets3 = [
    "1. Multi-Tier Query Normalization: Hierarchical regex (<0.1ms) + domain offence ontology.",
    "2. Deterministic Concordance Graph: Pure hash-table lookup for exact, split, merged, and repealed provisions.",
    "3. BM25 Statutory Retrieval (Design Choice): Eliminates dense embedding numerical blur across section numbers.",
    "4. Generative Answering: Strict [Act §Section] citation extraction grammar.",
    "5. Two-Layer Hard-Constraint Verifier:",
    "   • Layer 1: Closed-vocabulary statutory ID gating & repeal veto directives.",
    "   • Layer 1.5: Multi-citation cross-statute concordance consistency check.",
    "   • Layer 2: Penal duration & legal ingredient grounding bounding.",
    "   • Layer 2.5: Query-intent semantic relevance gating.",
    "6. Zero-Downtime Hot-Patch Refresh: Ingests new amendments dynamically in <5ms without re-indexing."
]
for b in bullets3:
    p = tf3.add_paragraph()
    p.text = b
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(8)

# ── Slide 4: Master Ablation Results ──────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
add_header(slide4, "Master Experimental Results (with 95% Wilson CIs)", "Rigorous Cross-Stage Empirical Evaluation")

# Add Table
rows_data = [
    ["Stage", "Configuration", "Testbed", "Accuracy / Metric", "95% Wilson CI", "Catch Rate", "FPR"],
    ["Stage 1", "Baseline LLM (Closed-Book)", "Dev Set (N=60)", "10.0% (6/60)", "[4.7% - 20.1%]", "N/A", "N/A"],
    ["Stage 2", "+BM25 RAG Context", "Dev Set (N=60)", "63.3% (38/60)", "[50.7% - 74.4%]", "N/A", "N/A"],
    ["Stage 3", "+Two-Layer Hard Verifier", "Stress Set (N=30)", "100.0% (30/30)", "[88.6% - 100.0%]", "100.0% (18/18)", "0.0% (0/12)"],
    ["Stage 4", "+Incremental Refresh", "Amendments (N=3)", "100.0% (3/3)", "Case Study", "100.0% (18/18)", "0.0% (0/12)"],
    ["Generalize", "CrPC <-> BNSS (Procedural)", "Procedural (N=30)", "100.0% (30/30)", "[88.6% - 100.0%]", "100.0%", "0.0% (0/30)"]
]

table_shape = slide4.shapes.add_table(len(rows_data), len(rows_data[0]), Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.2))
table = table_shape.table
for r_idx, row in enumerate(rows_data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                if r_idx == 0:
                    r.font.bold = True
                    r.font.size = Pt(12)
                    r.font.color.rgb = DARK_BLUE
                else:
                    r.font.size = Pt(11)
                    r.font.color.rgb = TEXT_DARK


tb4_sub = slide4.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
tf4_sub = tb4_sub.text_frame
tf4_sub.word_wrap = True
p = tf4_sub.paragraphs[0]
p.text = "• McNemar Paired Test (Stage 1 vs 2): chi2 = 28.26, p < 10^-6 (statistically significant jump).\n• Double-Blind Calibration: Cohen's Kappa kappa = 0.93 across N=20 calibrated test queries."
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(71, 85, 105)

# ── Slide 5: Live Demo & Viva Highlights ──────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
add_header(slide5, "Interactive Web UI Showcase (`app.py`) for Viva", "Real-Time Pipeline Inspection & Safety Demos")
tb5 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
tf5 = tb5.text_frame
tf5.word_wrap = True

bullets5 = [
    "• Live 5-Step Pipeline Inspection: Normalizer → Concordance → BM25 Retrieval → Generation → Verifier.",
    "• Sedition Repeal Interception (IPC §124A): Demonstrates active verifier veto and statutory advisory injection.",
    "• Split Section Ambiguity Breakdown: Demonstrates graded confidence on IPC §33 → BNS §2(1) & §2(25).",
    "• Zero-Downtime Amendment Toggle: Live switch between Base 2024 Gazette and Hot-Patched 2025 AI Amendment Index.",
    "• 67 Automated Unit Tests Passing (100% Pass Rate in 0.27s)."
]
for b in bullets5:
    p = tf5.add_paragraph()
    p.text = b
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(12)

out_pptx = "report/presentation_deck.pptx"
prs.save(out_pptx)
print(f"Presentation deck successfully updated: {out_pptx}")
